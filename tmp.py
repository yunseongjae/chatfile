import streamlit as st
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import PyPDFLoader, DirectoryLoader, UnstructuredPowerPointLoader
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from pdf2image import convert_from_path
from pptx import Presentation
from docx import Document
import tempfile
import fitz
from PIL import Image, ImageDraw
from pathlib import Path

DB_FAISS_PATH = "vectorstores/db_faiss"
embeddings = HuggingFaceEmbeddings(model_name = "jhgan/ko-sbert-sts", model_kwargs={'device': 'cpu'})

def create_vector_db():
    pdf_loader = DirectoryLoader(DATA_PATH, glob="*.pdf", loader_cls = PyPDFLoader, show_progress = True)
    ppt_loader = DirectoryLoader(DATA_PATH, glob="*.pptx", loader_cls = UnstructuredPowerPointLoader, show_progress = True)
    word_loader = DirectoryLoader(DATA_PATH, glob="*.docx", show_progress = True)

    loaders = [pdf_loader, ppt_loader, word_loader]
    
    documents = []
    for loader in loaders:
        documents.extend(loader.load())

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
    documents = text_splitter.split_documents(documents)
    
    db = FAISS.from_documents(documents, embeddings)
    db.save_local(DB_FAISS_PATH)



def query_embedding(query):
    if isinstance(query, list):
        query = query[0]

    embedded_query = embeddings.embed_query(query)
    return embedded_query

def extract_text_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_docx(docx_path):
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def show_pptx_docx_as_image(file_path):
    if file_path.endswith(".pptx"):
        text_to_show = extract_text_pptx(file_path)
    elif file_path.endswith(".docx"):
        text_to_show = extract_text_docx(file_path)
    st.text_area("원본", text_to_show)


def highlight_text(pdf_path, text_to_highlight):
    doc = fitz.open(pdf_path)
    for page in doc:
        text_instances = page.search_for(text_to_highlight)
        for inst in text_instances:
            highlight = page.add_highlight_annot(inst)
    pdf_path_highlighted = "highlighted.pdf"
    doc.save(pdf_path_highlighted)
    return pdf_path_highlighted

def show_highlighted_pdf(pdf_path, text_to_highlight, page_num):
    highlighted_pdf_path = highlight_text(pdf_path, text_to_highlight)
    images = convert_from_path(highlighted_pdf_path, first_page=page_num, last_page=page_num+1)
    st.image(images[0])
    st.image(images[1])

st.title("Jarvis DBQA")

folder_path = st.text_input("검색하고 싶은 폴더를 선택해주세요.")

if folder_path:
    DATA_PATH = folder_path

    if st.button('벡터 만들기'):
        create_vector_db()
        st.success('벡터 데이터베이스가 성공적으로 생성되었습니다.')

    db = FAISS.load_local(DB_FAISS_PATH, embeddings)
    query = st.text_input("검색하고 싶은 질문을 입력해주세요.")

    if query:
        embedded_query = query_embedding(query)
        docs = db.similarity_search(query)

        doc = docs[0]
        if 'page' in doc.metadata.keys():
            page_number = doc.metadata['page']
        file_path = doc.metadata['source']
        
        if file_path.endswith(".pdf"):
            text_to_highlight = doc.page_content
            st.write(doc.page_content)
            st.write(doc.metadata['source'])
            st.write(doc.metadata['page'])
            show_highlighted_pdf(file_path, text_to_highlight, page_number)
        else:
            st.text_area("유사문장: ",  doc.page_content)
            st.text_area("위치: ",  doc.metadata['source'])
            show_pptx_docx_as_image(file_path)
    








        
        
        # st.write(doc['page_content'])
